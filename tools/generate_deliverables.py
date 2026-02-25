import os
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
OUTPUT_PPTX = os.path.join(ROOT, "Parduzi_Secursign_Presentation.pptx")
OUTPUT_DOCX = os.path.join(ROOT, "Parduzi_Secursign_Rapport.docx")

SCREENSHOTS = [
    os.path.join(ROOT, "Capture d'écran 2026-02-24 153210.png"),
    os.path.join(ROOT, "Capture d'écran 2026-02-24 153339.png"),
    os.path.join(ROOT, "Capture d'écran 2026-02-24 153354.png"),
    os.path.join(ROOT, "Capture d'écran 2026-02-24 153410.png"),
]

TITLE = "PARDUZI SECURSIGN"
SUBTITLE = "Presentation de l'application Parduzi SecureSign"
AUTHOR = "Erion Parduzi"
CAMPUS = "Aix Ynov Campus"
DATE_STR = "27/02/2026"

CORP_BLUE = RGBColor(0x12, 0x2D, 0x4A)
ACCENT = RGBColor(0x0A, 0x91, 0xA9)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFB)


def build_presentation():
    prs = Presentation()
    prs.slide_width = PptInches(13.33)
    prs.slide_height = PptInches(7.5)

    def apply_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        apply_bg(slide, LIGHT_BG)
        shape = slide.shapes.add_shape(
            1, PptInches(0), PptInches(0), PptInches(13.33), PptInches(1.3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CORP_BLUE
        shape.line.fill.background()

        title_box = slide.shapes.add_textbox(PptInches(0.6), PptInches(0.2), PptInches(12), PptInches(1))
        title_tf = title_box.text_frame
        title_tf.text = TITLE
        title_tf.paragraphs[0].font.size = PptPt(36)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        sub_box = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.9), PptInches(12), PptInches(3))
        sub_tf = sub_box.text_frame
        sub_tf.text = SUBTITLE
        p = sub_tf.add_paragraph()
        p.text = f"{AUTHOR} - {CAMPUS}"
        p = sub_tf.add_paragraph()
        p.text = DATE_STR

    def add_section_title(text):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        apply_bg(slide, LIGHT_BG)
        bar = slide.shapes.add_shape(1, PptInches(0), PptInches(0), PptInches(13.33), PptInches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = CORP_BLUE
        bar.line.fill.background()

        box = slide.shapes.add_textbox(PptInches(0.6), PptInches(0.15), PptInches(12), PptInches(0.9))
        tf = box.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = PptPt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    def add_bullets(title_text, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        apply_bg(slide, LIGHT_BG)
        title = slide.shapes.add_textbox(PptInches(0.7), PptInches(0.5), PptInches(12), PptInches(0.7))
        tf = title.text_frame
        tf.text = title_text
        tf.paragraphs[0].font.size = PptPt(24)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = CORP_BLUE

        body = slide.shapes.add_textbox(PptInches(0.9), PptInches(1.5), PptInches(11.8), PptInches(5.5)).text_frame
        body.word_wrap = True
        for idx, item in enumerate(bullets):
            p = body.add_paragraph() if idx > 0 else body.paragraphs[0]
            p.text = item
            p.level = 0
            p.font.size = PptPt(18)

    def add_image_slide(title_text, image_path):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        apply_bg(slide, LIGHT_BG)
        title = slide.shapes.add_textbox(PptInches(0.7), PptInches(0.4), PptInches(12), PptInches(0.7))
        tf = title.text_frame
        tf.text = title_text
        tf.paragraphs[0].font.size = PptPt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = CORP_BLUE

        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, PptInches(0.8), PptInches(1.3), width=PptInches(11.8))
        else:
            add_bullets(title_text, ["Capture non trouvee : " + os.path.basename(image_path)])

    def add_code_slide(title_text, code_lines):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        apply_bg(slide, LIGHT_BG)
        title = slide.shapes.add_textbox(PptInches(0.7), PptInches(0.4), PptInches(12), PptInches(0.7))
        tf = title.text_frame
        tf.text = title_text
        tf.paragraphs[0].font.size = PptPt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = CORP_BLUE

        box = slide.shapes.add_textbox(PptInches(0.7), PptInches(1.4), PptInches(12), PptInches(5.6))
        code_tf = box.text_frame
        for i, line in enumerate(code_lines):
            p = code_tf.add_paragraph() if i > 0 else code_tf.paragraphs[0]
            p.text = line
            p.font.name = "Consolas"
            p.font.size = PptPt(14)

    add_title_slide()

    add_section_title("Contexte et objectifs")
    add_bullets("Problematique", [
        "Digitaliser le quitus de fin de chantier",
        "Fiabiliser les preuves (signature + GPS)",
        "Centraliser l'envoi et l'archivage",
        "Respecter le referentiel CDA"
    ])
    add_bullets("Objectifs", [
        "App mobile Android pour les techniciens",
        "Generation automatique de documents",
        "Envoi email et traçabilite",
        "Securite et authentification"
    ])

    add_section_title("Analyse de l'existant")
    add_bullets("Application web", [
        "Formulaire HTML + signatures",
        "Generation Excel",
        "Envoi par email",
        "Limites : pas mobile natif, UX terrain"
    ])

    add_section_title("Architecture globale")
    add_bullets("Composants", [
        "App Android (Kotlin + Compose)",
        "API Node.js (Express)",
        "Base PostgreSQL",
        "Services : mail, stockage fichiers"
    ])
    add_bullets("Flux", [
        "Saisie + signatures",
        "Appel API /api/generate",
        "Generation Excel + email",
        "Insertion BDD + hash de preuve"
    ])

    add_section_title("Conception applicative")
    add_bullets("UI/UX", [
        "Formulaire guide par sections",
        "Dropdown bailleurs",
        "Signature client et technicien",
        "Messages de statut et validations"
    ])
    add_bullets("Technique Android", [
        "Retrofit + OkHttp",
        "Compose Material 3",
        "GPS via Play Services",
        "Stockage local securise du token"
    ])

    add_section_title("Securite")
    add_bullets("Authentification", [
        "Inscription + verification email",
        "Mot de passe fort",
        "Hash bcrypt",
        "JWT et stockage chiffre",
        "Deep link reset mot de passe"
    ])
    add_bullets("Robustesse", [
        "Requetes SQL parametrees",
        "Hash de securite dans Excel",
        "Logs server et gestion erreurs"
    ])

    add_section_title("API et BDD")
    add_bullets("Endpoints", [
        "/api/generate, /api/auth/register, /api/auth/login",
        "/api/auth/forgot, /api/auth/reset, /api/auth/verify"
    ])
    add_bullets("Schema BDD", [
        "Table users (auth)",
        "Table interventions (metadonnees quitus)",
        "Hash et trace GPS"
    ])

    add_section_title("Demonstration")
    for idx, path in enumerate(SCREENSHOTS, start=1):
        add_image_slide(f"Capture {idx}", path)

    add_section_title("Extraits techniques")
    add_code_slide("Android - Envoi quitus", [
        "val payload = QuitusData(...) ",
        "val response = RetrofitClient.apiService.generateQuitus(payload)",
        "if (response.success) { ... }"
    ])
    add_code_slide("Node.js - Auth", [
        "const passwordHash = await bcrypt.hash(password, 12)",
        "await transporter.sendMail({ ... })",
        "return res.json({ success: true })"
    ])
    add_code_slide("Node.js - BDD", [
        "await pool.query('INSERT INTO interventions ...', values)",
        "const securityHash = crypto.createHash('sha256').update(raw).digest('hex')"
    ])

    add_section_title("Tests et validation")
    add_bullets("Verification", [
        "Build Android OK",
        "Appels API tests",
        "Verification email",
        "Reset mot de passe",
        "Generation Excel + email"
    ])

    add_section_title("Bilan et axes d'amelioration")
    add_bullets("Bilan", [
        "Application fonctionnelle terrain",
        "Traçabilite et securite renforces",
        "Process documentaire automatise"
    ])
    add_bullets("Axes d'amelioration", [
        "Migration iOS (KMP ou SwiftUI)",
        "Hebergement sur NAS securise",
        "Back-office web (validation, export, audit)",
        "Archivage et chiffrement long terme",
        "Statistiques et KPIs",
        "Mode hors-ligne + synchronisation"
    ])

    prs.save(OUTPUT_PPTX)


def add_doc_title(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(text)
    run.bold = True
    run.font.size = Pt(24)


def add_doc_subtitle(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(text)
    run.font.size = Pt(14)


def add_heading(doc, text, level=1):
    doc.add_heading(text, level=level)


def add_paragraph(doc, text):
    p = doc.add_paragraph(text)
    p.paragraph_format.space_after = Pt(10)


def add_code_block(doc, code_lines):
    for line in code_lines:
        p = doc.add_paragraph()
        run = p.add_run(line)
        run.font.name = 'Consolas'
        run._element.rPr.rFonts.set(qn('w:eastAsia'), 'Consolas')
        run.font.size = Pt(9)


def build_report():
    doc = Document()

    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style._element.rPr.rFonts.set(qn('w:eastAsia'), 'Calibri')
    style.font.size = Pt(11)

    add_doc_title(doc, "Rapport de projet")
    add_doc_subtitle(doc, TITLE)
    add_doc_subtitle(doc, f"{AUTHOR} - {CAMPUS} - {DATE_STR}")
    doc.add_page_break()

    sections = [
        ("1. Introduction", [
            "Ce rapport presente la realisation du projet Parduzi SecureSign, une application de quitus numerique visant a remplacer un processus papier par une solution mobile fiable et securisee.",
            "Le projet s'inscrit dans le cadre du titre RNCP niveau 6 (CDA). Il couvre l'analyse, la conception, le developpement, la securisation, les tests et la mise en production d'une solution complete.",
            "L'approche retenue met l'accent sur la traçabilite des interventions, la simplification des demarches terrain et la fiabilisation des preuves de fin de chantier."
        ]),
        ("2. Contexte et besoins", [
            "Les equipes terrain doivent valider des interventions de fin de chantier par un quitus signe. L'ancienne solution etait un formulaire web, peu adapte a l'usage mobile.",
            "Les besoins principaux sont : saisie rapide des donnees, signatures client et technicien, envoi automatique par email et archivage des preuves.",
            "Des contraintes fortes existent : utilisation en mobilite, reseau parfois instable, et besoin d'une preuve inviolable."
        ]),
        ("3. Objectifs du projet", [
            "1) Concevoir une application mobile Android intuitive.",
            "2) Automatiser la generation de documents.",
            "3) Securiser les acces et la transmission.",
            "4) Assurer la conformite au referentiel CDA.",
            "5) Centraliser les donnees et faciliter l'audit."
        ]),
        ("4. Analyse de l'existant", [
            "L'application web initiale permettait la saisie et la generation d'un fichier Excel, mais la signature et l'usage mobile restaient limites.",
            "Les principales limites constataes : ergonomie terrain insuffisante, pas de stockage securise de session, et integration mobile incomplète.",
            "L'objectif a ete de porter les fonctionnalites essentielles en natif Android pour un usage terrain fiable."
        ]),
        ("5. Architecture globale", [
            "Le systeme est compose d'une app Android, d'une API Node.js (Express) et d'une base PostgreSQL.",
            "L'application envoie les donnees, l'API genere les documents et conserve une trace securisee en base.",
            "Le serveur gere l'emailing, l'archivage et la generation de la preuve."
        ]),
        ("6. Choix techniques", [
            "Android : Kotlin + Jetpack Compose pour l'UI, Material 3 pour l'ergonomie.",
            "Backend : Node.js + Express pour la rapidite de developpement et l'integration email.",
            "BDD : PostgreSQL pour la robustesse, la conformite SQL et la fiabilite.",
            "Reseau : Retrofit + OkHttp avec gestion des timeouts."
        ]),
        ("7. Conception de l'application mobile", [
            "L'interface est decoupee en sections claires : informations dossier, locataire, details et signatures.",
            "Les composants Compose sont reutilisables (InputGroup, CustomTextField, SignatureModal).",
            "La validation locale bloque l'envoi si les champs obligatoires ou les signatures manquent.",
            "Le workflow vise la rapidite et la reduction des erreurs de saisie."
        ]),
        ("8. Capture de signatures", [
            "La capture de signature se fait via un canvas Compose. Les traits sont convertis en base64.",
            "Les signatures client et technicien sont stockees separement, puis injectees dans le document final.",
            "Cette approche garantit une preuve visuelle et une traçabilite directe."
        ]),
        ("9. Envoi du quitus", [
            "Une fois le formulaire complet, l'application construit un payload JSON et appelle l'API /api/generate.",
            "L'API genere le fichier Excel, envoie l'email et retourne un statut d'execution.",
            "Un message utilisateur signale la reussite ou l'erreur."
        ]),
        ("10. Backend Node.js", [
            "Le serveur Express gere la creation des quitus, l'injection des donnees dans un fichier Excel et l'envoi par email.",
            "Une empreinte SHA-256 est ajoutee pour garantir l'integrite des documents.",
            "Le serveur controle les erreurs, trace les exceptions et conserve les fichiers generes."
        ]),
        ("11. Base de donnees", [
            "PostgreSQL stocke les interventions et les metadonnees : bailleur, client, adresse, GPS, chemin du fichier et hash securite.",
            "La table users stocke l'authentification avec hash de mot de passe et verification email.",
            "Les requetes parametrees evitent les injections SQL."
        ]),
        ("12. Authentification et securite", [
            "L'inscription impose un mot de passe fort, stocke par hash bcrypt.",
            "Une verification email active le compte avant la premiere connexion.",
            "L'authentification repose sur JWT et le token est stocke dans un coffre chiffre Android.",
            "Le flux 'mot de passe oublie' utilise un token a usage unique et expire."
        ]),
        ("13. Deconnexion", [
            "Un bouton de deconnexion est disponible dans l'application.",
            "Il supprime le token local et renvoie vers l'ecran d'authentification.",
            "Cela renforce la securite en cas de perte ou de partage du terminal."
        ]),
        ("14. API et endpoints", [
            "POST /api/generate : generation du quitus et envoi email.",
            "POST /api/auth/register : inscription et envoi email de verification.",
            "POST /api/auth/login : connexion et generation JWT.",
            "POST /api/auth/forgot : lien de reinitialisation.",
            "POST /api/auth/reset : changement du mot de passe.",
            "GET /api/auth/verify : activation du compte."
        ]),
        ("15. Qualite et conformite", [
            "Le projet respecte les bonnes pratiques CDA : separation des couches, validation, securisation des donnees.",
            "La documentation accompagne l'installation, les tests et l'usage terrain.",
            "Les choix techniques privilegient la maintenabilite et la fiabilite."
        ]),
        ("16. Tests et validation", [
            "Tests manuels sur le cycle complet : saisie, signatures, envoi email, generation Excel.",
            "Tests d'authentification : register, login, forgot, reset, verification.",
            "Tests de robustesse reseau : timeouts et erreurs traitees proprement."
        ]),
        ("17. Resultats et valeur ajoutee", [
            "Application mobile operationnelle pour le terrain.",
            "Process documentaire automatise et fiable.",
            "Traçabilite renforcée par la combinaison signature + GPS + hash."
        ]),
        ("18. Axes d'amelioration", [
            "Migration iOS (KMP ou SwiftUI) pour couvrir tous les appareils.",
            "Hebergement sur NAS securise pour l'archivage interne.",
            "Back-office web pour le suivi des interventions.",
            "Mode hors-ligne avec synchronisation.",
            "Tableau de bord KPI et statistiques."
        ])
    ]

    for title, paras in sections:
        add_heading(doc, title, level=1)
        for para in paras:
            add_paragraph(doc, para)
        doc.add_page_break()

    add_heading(doc, "19. Extraits de code", level=1)
    add_paragraph(doc, "Extrait Android - envoi de quitus")
    add_code_block(doc, [
        "val payload = QuitusData(...) ",
        "val response = RetrofitClient.apiService.generateQuitus(payload)",
        "if (response.success) { /* succes */ }"
    ])

    add_paragraph(doc, "Extrait Node.js - inscription")
    add_code_block(doc, [
        "const passwordHash = await bcrypt.hash(password, 12)",
        "await transporter.sendMail({ ... })",
        "return res.json({ success: true })"
    ])

    add_paragraph(doc, "Extrait Node.js - insertion BDD")
    add_code_block(doc, [
        "await pool.query('INSERT INTO interventions (...) VALUES (...)', values)",
        "const securityHash = crypto.createHash('sha256').update(raw).digest('hex')"
    ])

    doc.add_page_break()

    add_heading(doc, "20. Annexes - Captures", level=1)
    for path in SCREENSHOTS:
        if os.path.exists(path):
            add_heading(doc, f"Capture - {os.path.basename(path)}", level=2)
            doc.add_picture(path, width=Inches(6.5))
            doc.add_paragraph("")

    doc.save(OUTPUT_DOCX)


if __name__ == "__main__":
    build_presentation()
    build_report()
    print("Generated:")
    print(OUTPUT_PPTX)
    print(OUTPUT_DOCX)
